from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from transformers import pipeline
from fpdf import FPDF

import docx2txt
import os
import uuid
import fitz  # PyMuPDF for PDF text extraction
import markdown
import pptx
from io import BytesIO

app = FastAPI()

# Allow frontend access - configured to work with Next.js frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],  # Add your production domain when deployed
    allow_methods=["*"],
    allow_headers=["*"],
)

# Load summarization model (offline)
summarizer = pipeline("summarization", model="facebook/bart-large-cnn", device=-1)

# ---------------------------- Utilities ---------------------------- #

def chunk_text(text, max_tokens=900):
    sentences = text.split('. ')
    chunks, chunk = [], ""
    for s in sentences:
        if len(chunk.split()) + len(s.split()) < max_tokens:
            chunk += s + ". "
        else:
            chunks.append(chunk.strip())
            chunk = s + ". "
    chunks.append(chunk.strip())
    return chunks

def summarize_chunk(chunk):
    result = summarizer(chunk, max_length=512, min_length=100, do_sample=False)[0]['summary_text']
    return result.strip()

def extract_text_from_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_ppt(file_path: str) -> str:
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def extract_text(file_path: str, extension: str) -> str:
    if extension == ".docx":
        return docx2txt.process(file_path)
    elif extension == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension == ".pptx" or extension == ".ppt":
        return extract_text_from_ppt(file_path)
    elif extension == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    elif extension == ".md":
        with open(file_path, "r", encoding="utf-8") as f:
            return markdown.markdown(f.read())
    else:
        raise ValueError("Unsupported file type")

def create_pdf(summary_text: str, style: str) -> str:
    pdf = FPDF()
    pdf.add_page()

    if style == "minimal":
        pdf.set_font("Arial", size=11)
        pdf.set_text_color(0, 0, 0)
        line_height = 8

    elif style == "aesthetic":
        pdf.set_font("Helvetica", style='B', size=12)
        pdf.set_text_color(60, 60, 60)
        line_height = 9

    elif style == "elegant":
        pdf.set_font("Times", style='I', size=12)
        pdf.set_text_color(40, 40, 80)
        line_height = 10

    elif style == "academic":
        pdf.set_font("Courier", size=11)
        pdf.set_text_color(20, 20, 20)
        line_height = 7
        pdf.set_left_margin(15)
        pdf.set_right_margin(15)
        # Add academic title
        pdf.set_font("Courier", 'B', 14)
        pdf.cell(0, 10, "Academic Summary", ln=True, align='C')
        pdf.ln(5)
        pdf.set_font("Courier", size=11)  # Reset font for content

    else:
        pdf.set_font("Arial", size=11)
        line_height = 8

    for line in summary_text.split('\n'):
        pdf.multi_cell(0, line_height, txt=line)

    os.makedirs("outputs", exist_ok=True)
    output_path = f"./outputs/summary_{uuid.uuid4().hex[:8]}_{style}.pdf"
    pdf.output(output_path)
    return output_path

# Format summary into a structured format for the frontend
def format_summary_for_frontend(full_summary: str):
    parts = full_summary.split("\n\n📌 Part")
    formatted_summary = []
    
    if len(parts) > 1:
        # Introduction part
        if parts[0].strip():
            formatted_summary.append({
                "id": 0,
                "title": "Overview",
                "points": [p.strip() for p in parts[0].strip().split("\n") if p.strip()]
            })
            
        # Process other parts
        for i, part in enumerate(parts[1:], 1):
            part_text = part.strip()
            if part_text:
                # Split at the first colon to separate the title
                part_content = part_text.split(":", 1)
                
                title = f"Section {i}"
                content = part_text
                
                if len(part_content) > 1:
                    # Try to extract a better title from the first line
                    title = part_content[0].replace(f"{i}:", "").strip()
                    content = part_content[1].strip()
                
                # Split content into bullet points for better readability
                points = [line.strip() for line in content.split(". ") if line.strip()]
                
                formatted_summary.append({
                    "id": i,
                    "title": title,
                    "points": points
                })
    else:
        # No parts found, use the whole summary
        formatted_summary.append({
            "id": 1,
            "title": "Document Summary",
            "points": [p.strip() for p in full_summary.split("\n") if p.strip()]
        })
            
    return formatted_summary


@app.post("/summarize/")
async def summarize_file(
    file: UploadFile = File(...),
    style: str = Form("minimal")  # Choose: minimal, aesthetic, elegant, academic
):
    try:
        ext = os.path.splitext(file.filename)[1].lower()
        allowed_exts = [".docx", ".pdf", ".pptx", ".ppt", ".txt", ".md"]

        if ext not in allowed_exts:
            return JSONResponse(status_code=400, content={"error": "Unsupported file format"})

        os.makedirs("uploads", exist_ok=True)
        file_path = f"./uploads/{uuid.uuid4().hex[:6]}_{file.filename}"
        with open(file_path, "wb") as f:
            f.write(await file.read())

        raw_text = extract_text(file_path, ext)
        if not raw_text.strip():
            return JSONResponse(status_code=400, content={"error": "File contains no readable text"})

        chunks = chunk_text(raw_text)
        full_summary = ""

        for i, chunk in enumerate(chunks):
            part_summary = summarize_chunk(chunk)
            full_summary += f"\n\n📌 Part {i+1}:\n{part_summary}"

        # Create formatted summary for frontend
        formatted_summary = format_summary_for_frontend(full_summary)
        
        # Generate PDF
        pdf_path = create_pdf(full_summary, style)
        
        return {
            "summary": formatted_summary,
            "pdf_url": f"/api/download/{os.path.basename(pdf_path)}"
        }

    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.get("/download/{filename}")
async def download_file(filename: str):
    file_path = f"./outputs/{filename}"
    return FileResponse(file_path, media_type='application/pdf', filename=filename)

# Health check endpoint
@app.get("/health")
async def health_check():
    return {"status": "ok", "timestamp": "2025-04-06 13:56:17"}